import io
import calendar
import datetime

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from app import models
from app.database import get_db
from app.routes.annual_plan import _latest_week, _kpi_to_out

router = APIRouter(prefix="/api/export", tags=["export"])

# ─── Color palette (light mode) ──────────────────────────────────────────────

KPI_COLORS: dict[int, RGBColor] = {
    1: RGBColor(0x22, 0xD3, 0xEE),
    2: RGBColor(0x38, 0xBD, 0xF8),
    3: RGBColor(0x60, 0xA5, 0xFA),
    4: RGBColor(0x81, 0x8C, 0xF8),
    5: RGBColor(0xA7, 0x8B, 0xFA),
}

WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
PANEL_BG   = RGBColor(0xF8, 0xFA, 0xFC)
PANEL_HDR  = RGBColor(0xF1, 0xF5, 0xF9)
BORDER_CLR = RGBColor(0xE2, 0xE8, 0xF0)
GRID_CLR   = RGBColor(0xE2, 0xE8, 0xF0)
MONTH_A    = RGBColor(0xF8, 0xFA, 0xFC)
MONTH_B    = RGBColor(0xF1, 0xF5, 0xF9)
ROW_A      = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_PRI   = RGBColor(0x0F, 0x17, 0x2A)
TEXT_SEC   = RGBColor(0x47, 0x56, 0x69)
TEXT_MUT   = RGBColor(0x94, 0xA3, 0xB8)
TODAY_CLR  = RGBColor(0xEF, 0x44, 0x44)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
YELLOW     = RGBColor(0xEA, 0xB3, 0x08)
EMERALD    = RGBColor(0x10, 0xB9, 0x81)
BLUE_FG    = RGBColor(0x25, 0x63, 0xEB)
STAR_CLR   = RGBColor(0xF5, 0x9E, 0x0B)
NUM_BG     = RGBColor(0xDB, 0xEA, 0xFE)
NUM_FG     = RGBColor(0x25, 0x63, 0xEB)

STATUS_CLR   = {'in_progress': RGBColor(0x3B, 0x82, 0xF6),
                'completed':   EMERALD, 'not_started': TEXT_MUT}
STATUS_BG    = {'in_progress': RGBColor(0xEF, 0xF6, 0xFF),
                'completed':   RGBColor(0xEC, 0xFD, 0xF5), 'not_started': PANEL_BG}
STATUS_LABEL = {'in_progress': '進行中', 'completed': '已完成', 'not_started': '未開始'}

MONTH_NAMES = ['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec']
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
ML = Inches(0.30)

# ─── Geometry ────────────────────────────────────────────────────────────────

def _is_leap(y): return (y % 4 == 0 and y % 100 != 0) or y % 400 == 0
def _days(y):    return 366 if _is_leap(y) else 365

def _date_frac(d: datetime.date, year: int) -> float:
    return max(0.0, min(1.0, (d - datetime.date(year, 1, 1)).days / _days(year)))

def _month_frac(m, year): return _date_frac(datetime.date(year, m, 1), year)
def _month_w_frac(m, year): return calendar.monthrange(year, m)[1] / _days(year)

def _seg_bar(s, e, year, tl_x, tl_w_in):
    if s is None or e is None:
        return None, None
    yr0, yr1 = datetime.date(year, 1, 1), datetime.date(year, 12, 31)
    cs, ce = max(s, yr0), min(e, yr1)
    if cs > ce:
        return None, None
    sf = _date_frac(cs, year)
    ef = _date_frac(ce + datetime.timedelta(days=1), year)
    bx = int(tl_x + Inches(tl_w_in * sf))
    bw = max(int(Inches(tl_w_in * (ef - sf))), 50_000)
    return bx, bw

def _pct_color(pct):
    if not pct:   return TEXT_MUT
    if pct >= 80: return EMERALD
    if pct >= 50: return YELLOW
    return AMBER

# ─── Drawing primitives ──────────────────────────────────────────────────────

def _rect(slide, x, y, w, h, fill, *, border=None, bw=0.5):
    s = slide.shapes.add_shape(1, int(x), int(y), max(int(w), 9000), max(int(h), 9000))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(bw)
    else:       s.line.fill.background()
    return s

def _oval_shape(slide, x, y, d, fill, *, border=None, bw=2.5):
    d = max(int(d), 9000)
    s = slide.shapes.add_shape(9, int(x), int(y), d, d)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(bw)
    else:       s.line.fill.background()
    return s

def _text(slide, text, x, y, w, h, *, size=9, bold=False,
          color=TEXT_PRI, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(int(x), int(y), max(int(w), 9000), max(int(h), 9000))
    tf  = txb.text_frame
    tf.word_wrap  = wrap
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top  = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

def _shape_text(shape, text, size=9, bold=False, color=TEXT_PRI, align=PP_ALIGN.CENTER):
    """Write vertically-centered text into an existing shape."""
    tf = shape.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.margin_left = tf.margin_right = Pt(2)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

def _circle_label(slide, x, y, d, fill, text, text_color, *,
                  border=None, bw=0, size=8, bold=True):
    """Filled or bordered circle with centered label."""
    s = _oval_shape(slide, x, y, d, fill, border=border, bw=bw)
    _shape_text(s, text, size=size, bold=bold, color=text_color)

def _donut(slide, x, y, d, pct, color=None):
    """White circle, colored border, % text inside (donut style)."""
    c = color if color else _pct_color(pct)
    t = f"{pct}%" if pct is not None else "—"
    s = _oval_shape(slide, x, y, d, WHITE, border=c, bw=2.5)
    _shape_text(s, t, size=7, bold=True, color=c)

def _badge_rect(slide, x, y, w, h, fill, text, text_color, *,
                border=None, bw=0.5, size=8.5, bold=False):
    """Rounded rect with centered text (status / avg badges)."""
    s = slide.shapes.add_shape(5, int(x), int(y), max(int(w), 9000), max(int(h), 9000))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = Pt(bw)
    else:       s.line.fill.background()
    _shape_text(s, text, size=size, bold=bold, color=text_color)

def _light_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide

# ─── Gantt helpers ───────────────────────────────────────────────────────────

def _month_headers(slide, year, tl_x, tl_w_in, y, h):
    for m in range(1, 13):
        mx = int(tl_x + Inches(tl_w_in * _month_frac(m, year)))
        mw = int(Inches(tl_w_in * _month_w_frac(m, year)))
        _rect(slide, mx, y, mw, h, MONTH_A if m % 2 else MONTH_B)
        _text(slide, MONTH_NAMES[m - 1], mx, y, mw, h,
              size=6.5, color=TEXT_MUT, align=PP_ALIGN.CENTER)

def _grid_lines(slide, year, tl_x, tl_w_in, y_top, height):
    for m in range(2, 13):
        gx = int(tl_x + Inches(tl_w_in * _month_frac(m, year)))
        _rect(slide, gx, y_top, 9000, int(height), GRID_CLR)

def _today_line(slide, year, tl_x, tl_w_in, y_top, height):
    today = datetime.date.today()
    if today.year != year:
        return
    tx = int(tl_x + Inches(tl_w_in * _date_frac(today, year)))
    _rect(slide, tx, y_top, 18000, int(height), TODAY_CLR)

# ─── Slide 1: Gantt (light) ──────────────────────────────────────────────────

def _build_gantt(prs, year: int, groups):
    slide = _light_slide(prs)

    _text(slide, "年度計劃", ML, Inches(0.10), Inches(5), Inches(0.55),
          size=24, bold=True, color=TEXT_PRI)

    KPI_W_IN  = 1.40
    TASK_W_IN = 1.65
    PCT_W_IN  = 0.55
    TL_W_IN   = 13.33 - 0.30 - KPI_W_IN - TASK_W_IN - PCT_W_IN - 0.30  # ≈ 9.13
    TASK_X = int(ML + Inches(KPI_W_IN))
    PCT_X  = int(TASK_X + Inches(TASK_W_IN))
    TL_X   = int(PCT_X + Inches(PCT_W_IN))
    KPI_W  = Inches(KPI_W_IN)
    TASK_W = Inches(TASK_W_IN)
    PCT_W  = Inches(PCT_W_IN)
    FULL_W = Inches(KPI_W_IN + TASK_W_IN + PCT_W_IN + TL_W_IN)

    HDR_Y = Inches(0.72)
    HDR_H = Inches(0.28)

    # Column header row
    _rect(slide, ML, HDR_Y, FULL_W, HDR_H, PANEL_HDR)
    _text(slide, "KPI",   ML,      HDR_Y, KPI_W,  HDR_H, size=6.5, bold=True, color=TEXT_MUT)
    _text(slide, "TASK",  TASK_X,  HDR_Y, TASK_W, HDR_H, size=6.5, bold=True, color=TEXT_MUT)
    _text(slide, "工量%",  PCT_X,   HDR_Y, PCT_W,  HDR_H, size=6.5, bold=True, color=TEXT_MUT, align=PP_ALIGN.CENTER)
    _month_headers(slide, year, TL_X, TL_W_IN, HDR_Y, HDR_H)

    # Row list
    all_rows: list[tuple] = []
    for group, color in groups:
        all_rows.append(('kpi', group, None, None, color))
        for task in group.tasks:
            dated = [i for i in task.items if i.segments]
            if not dated:
                continue
            all_rows.append(('task', group, task, None, color))
            for item in dated:
                all_rows.append(('item', group, task, item, color))

    if not all_rows:
        return

    CONTENT_Y = float(HDR_Y + HDR_H)
    row_h = max(int((Inches(7.42) - CONTENT_Y) / len(all_rows)), int(Inches(0.18)))

    for ri, (kind, group, task, item, color) in enumerate(all_rows):
        y  = int(CONTENT_Y) + row_h * ri
        rh = row_h
        row_bg = ROW_A if ri % 2 == 0 else WHITE
        _rect(slide, ML, y, FULL_W, rh, row_bg)
        _grid_lines(slide, year, TL_X, TL_W_IN, y, rh)

        if kind == 'kpi':
            _text(slide, group.kpi_title.split('. ', 1)[-1],
                  int(ML + Inches(0.06)), y, int(KPI_W - Inches(0.08)), rh,
                  size=8, bold=True, color=color)
            if group.percentage is not None:
                _text(slide, f"{group.percentage}%", PCT_X, y, PCT_W, rh,
                      size=8, bold=True, color=color, align=PP_ALIGN.CENTER)

        elif kind == 'task':
            _text(slide, task.title, int(TASK_X + Inches(0.05)), y,
                  int(TASK_W - Inches(0.07)), rh, size=7, color=TEXT_SEC)

        else:
            bar_y = y + int(rh * 0.10)
            bar_h = int(rh * 0.80)
            for seg in item.segments:
                bx, bw = _seg_bar(seg.start_date, seg.end_date, year, TL_X, TL_W_IN)
                if bx is None:
                    continue
                _rect(slide, bx, bar_y, bw, bar_h, color)
                if bw > Inches(0.55):
                    _text(slide, item.content, int(bx + Inches(0.05)), bar_y,
                          int(bw - Inches(0.07)), bar_h,
                          size=6.5, color=WHITE, wrap=False)

        _today_line(slide, year, TL_X, TL_W_IN, y, rh)

# ─── Slides 2-6: KPI Detail (light) ─────────────────────────────────────────

def _build_kpi_detail(prs, kpi: models.KPI, kpi_num: int):
    slide = _light_slide(prs)
    color  = KPI_COLORS.get(kpi_num, RGBColor(0x60, 0xA5, 0xFA))
    status = kpi.status or 'not_started'

    # ─ Title ─────────────────────────────────────────────────────────────────
    _text(slide, kpi.title, ML, Inches(0.12), Inches(11), Inches(0.58),
          size=22, bold=True, color=TEXT_PRI)

    # ─ Status badge ──────────────────────────────────────────────────────────
    s_color = STATUS_CLR.get(status, TEXT_MUT)
    s_bg    = STATUS_BG.get(status, PANEL_BG)
    s_label = STATUS_LABEL.get(status, '未開始')
    _badge_rect(slide, ML, Inches(0.78), Inches(1.20), Inches(0.30),
                s_bg, f"● {s_label}", s_color, border=s_color, bw=0.5, size=8.5)

    # ─ Average % from highlights ──────────────────────────────────────────────
    pct_vals = [hl.percentage for hl in kpi.highlights if hl.percentage is not None]
    avg_pct  = round(sum(pct_vals) / len(pct_vals)) if pct_vals else None

    # ─ Panel layout (dynamic height, vertically centred) ──────────────────────
    PHDR_H = Inches(0.46)
    SKPI_H = Inches(0.38)
    ITEM_H = Inches(0.30)
    HL_H   = Inches(0.50)
    HL_GAP = Inches(0.04)

    left_h  = sum(float(SKPI_H) + len(sk.items) * float(ITEM_H) for sk in kpi.sub_kpis)
    right_h = len(kpi.highlights) * (float(HL_H) + float(HL_GAP))
    needed  = float(PHDR_H) + max(left_h, right_h) + float(Inches(0.24))

    PANEL_H   = min(max(needed, float(Inches(2.50))), float(Inches(5.20)))
    CONTENT_START = float(Inches(1.15))
    CONTENT_END   = float(SLIDE_H) - float(Inches(0.35))
    PANEL_TOP = CONTENT_START + (CONTENT_END - CONTENT_START - PANEL_H) / 2
    PANEL_TOP = max(PANEL_TOP, CONTENT_START)
    LEFT_X    = ML
    LEFT_W    = Inches(6.08)
    GAP       = Inches(0.25)
    RIGHT_X   = LEFT_X + LEFT_W + GAP
    RIGHT_W   = SLIDE_W - RIGHT_X - ML

    # Panel backgrounds
    _rect(slide, LEFT_X,  PANEL_TOP, LEFT_W,  PANEL_H, PANEL_BG, border=BORDER_CLR)
    _rect(slide, RIGHT_X, PANEL_TOP, RIGHT_W, PANEL_H, PANEL_BG, border=BORDER_CLR)

    # ─ Left panel header: KPI 指標 ───────────────────────────────────────────
    _rect(slide, LEFT_X, PANEL_TOP, LEFT_W, PHDR_H, PANEL_HDR)
    # Bottom separator line
    _rect(slide, LEFT_X, int(PANEL_TOP + PHDR_H - Inches(0.01)), LEFT_W, Inches(0.01), BORDER_CLR)

    _text(slide, "⊙  KPI 指標", int(LEFT_X + Inches(0.14)), PANEL_TOP,
          Inches(3.0), PHDR_H, size=9.5, bold=True, color=BLUE_FG)

    # "平均" label + donut (average of highlight percentages)
    if avg_pct is not None:
        pct_color = _pct_color(avg_pct)
        D = Inches(0.40)
        d_x = int(LEFT_X + LEFT_W - D - Inches(0.12))
        d_y = int(PANEL_TOP + (PHDR_H - D) / 2)
        _donut(slide, d_x, d_y, D, avg_pct, pct_color)
        _text(slide, "平均", int(d_x - Inches(0.55)), int(PANEL_TOP),
              Inches(0.50), PHDR_H, size=8, color=TEXT_MUT, align=PP_ALIGN.RIGHT)

    # ─ Left panel content ────────────────────────────────────────────────────
    row_y  = PANEL_TOP + float(PHDR_H) + float(Inches(0.06))
    max_y  = PANEL_TOP + PANEL_H - float(Inches(0.06))
    SKPI_H = Inches(0.38)
    ITEM_H = Inches(0.30)
    PAD_L  = Inches(0.14)

    for sub_kpi in kpi.sub_kpis:
        if row_y + float(SKPI_H) > max_y:
            break
        _text(slide, f"{sub_kpi.sub_id}  {sub_kpi.title}",
              int(LEFT_X + PAD_L), int(row_y),
              int(LEFT_W - PAD_L - Inches(0.08)), int(SKPI_H),
              size=9, bold=True, color=TEXT_PRI)
        row_y += float(SKPI_H)

        for item in sub_kpi.items:
            if row_y + float(ITEM_H) > max_y:
                break
            _text(slide, f"  >  {item.content}",
                  int(LEFT_X + PAD_L), int(row_y),
                  int(LEFT_W - PAD_L - Inches(0.08)), int(ITEM_H),
                  size=8, color=TEXT_SEC)
            row_y += float(ITEM_H)

    # ─ Right panel header: HIGHLIGHT ─────────────────────────────────────────
    _rect(slide, RIGHT_X, PANEL_TOP, RIGHT_W, PHDR_H, PANEL_HDR)
    _rect(slide, RIGHT_X, int(PANEL_TOP + PHDR_H - Inches(0.01)), RIGHT_W, Inches(0.01), BORDER_CLR)
    _text(slide, "★  HIGHLIGHT", int(RIGHT_X + Inches(0.14)), PANEL_TOP,
          Inches(4.0), PHDR_H, size=9.5, bold=True, color=STAR_CLR)

    # ─ Right panel content ───────────────────────────────────────────────────
    NC_D   = Inches(0.30)   # number circle diameter
    DOT_D  = Inches(0.38)   # percentage donut diameter
    row_y  = PANEL_TOP + float(PHDR_H) + float(Inches(0.06))
    max_y  = PANEL_TOP + PANEL_H - float(Inches(0.06))

    for i, hl in enumerate(kpi.highlights, 1):
        if row_y + float(HL_H) > max_y:
            break
        rh = int(HL_H)

        row_bg = ROW_A if i % 2 == 1 else PANEL_BG
        _rect(slide, RIGHT_X, int(row_y), RIGHT_W, rh, row_bg)

        # Number circle (light blue)
        nc_x = int(RIGHT_X + Inches(0.12))
        nc_y = int(row_y) + int((rh - int(NC_D)) / 2)
        _circle_label(slide, nc_x, nc_y, NC_D, NUM_BG, str(i), NUM_FG, size=8)

        # Content text
        text_x = nc_x + int(NC_D) + int(Inches(0.10))
        text_w = int(float(RIGHT_W) - float(NC_D) - float(DOT_D) - Inches(0.52))
        _text(slide, hl.content, text_x, int(row_y), text_w, rh, size=8, color=TEXT_PRI)

        # Percentage donut
        if hl.percentage is not None:
            dot_x = int(RIGHT_X + float(RIGHT_W) - float(DOT_D) - Inches(0.12))
            dot_y = int(row_y) + int((rh - int(DOT_D)) / 2)
            _donut(slide, dot_x, dot_y, DOT_D, hl.percentage)

        row_y += float(HL_H) + Inches(0.04)

# ─── Endpoint ─────────────────────────────────────────────────────────────────

@router.get("/annual-plan/{year}/pptx")
def export_annual_plan_pptx(year: int, db: Session = Depends(get_db)):
    week = _latest_week(year, db)
    actual_year = year
    if not week:
        week = _latest_week(year - 1, db)
        actual_year = year - 1
    if not week:
        raise HTTPException(status_code=404, detail="No annual plan data found")

    kpis: list[models.KPI] = (
        db.query(models.KPI)
        .filter(models.KPI.week_id == week.id)
        .order_by(models.KPI.number)
        .all()
    )
    annual_groups = [_kpi_to_out(k, week.id) for k in kpis]

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    pairs = [(g, KPI_COLORS.get(g.kpi_number, RGBColor(0x60, 0xA5, 0xFA)))
             for g in annual_groups]
    _build_gantt(prs, actual_year, pairs)

    for kpi in kpis:
        _build_kpi_detail(prs, kpi, kpi.number)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    filename = f"annual-plan-{actual_year}.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
